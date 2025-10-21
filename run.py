from pptx import Presentation, parts
from pptx.enum.text import PP_ALIGN

from pptx.util import Pt

from openpyxl import load_workbook
import copy
import os 
import math

import requests
from bs4 import BeautifulSoup
import re

## Prepare data
## Our xlsx file has the following format
## | No | Name | Affiliation | ETC | ETC | ... 
"""
<thead>
        <tr>
            <th>#</th>
            <th>Name</th>
            <th>Job Title</th>
            <th>E-Mail</th>
            <th>Affiliation</th>
            <th>Registration Type</th>
        </tr>
    </thead>
    <tbody id='personal_infos_tbody'><tr>
            <td>1</td> 
            <td>Jonghyeok Park</td>
            <td>Assistant Professor</td>
            <td>jonghyeok_park@korea.ac.kr</td>
            <td>Korea University</td>
            <td>FULL</td>
"""

def check_korean_name(name):
    return all('\uac00' <= char <= '\ud7a3' or char.isspace() for char in name)

def swap_korean_name(name):
    token = name.split()
    if len(token) == 2:  # convert first name + last name to last name + first name
        return f"{token[1]} {token[0]}"  # 이름 성으로 변경
    return name

names = []
affiliations = []

def prepare_data_v1():
    total = 0
    url = "https://sample-url.php"

    '''
    아래와 같은 등록자 명단 관리한다고 가정함.
    <thead>
        <tr>
            <th>Select</th>
            <th>#</th>
            <th>Name</th>
            <th>Korean Name</th>
            <th>Job Title</th>
            <th>E-Mail</th>
            <th>Affiliation</th>
            <th>Registration Type</th>
        </tr>
    </thead>
    <tbody id='personal_infos_tbody'>
    <tr>
                    <td><input type='checkbox' name='ids[]' value='1'></td>
                    <td>1</td>
                    <td>Jonghyeok Park</td>
                    <td>박종혁</td>
                    <td>조교수</td>
                    <td>jonghyeok_park@korea.ac.kr</td>
                    <td>Korea University</td>
                    <td>FULL</td>
    </tr> ...
    '''


    response = requests.get(url)
    if response.status_code != 200:
        print(f"Failed to retrieve the server. Status code: {response.status_code}")
        return None

    soup = BeautifulSoup(response.content, 'html.parser')
    tbl_body = soup.find('tbody', {'id': 'personal_infos_tbody'})
    tbl_rows = tbl_body.find_all('tr')

    for row in tbl_rows:
        cells = row.find_all('td')
        if len(cells) >= 5:  # Name과 Affiliation이 포함된 5개 이상의 열이 있는 경우
            name = cells[2].get_text(strip=True)  # cell[1] means name
            affiliation = cells[6].get_text(strip=True)  # dell[4] means affiliation

            #if check_korean_name(name):
            #    name = swap_korean_name(name)

            names.append(name)
            affiliations.append(affiliation)
            total+=1
    # debug
    #print(f"names:{names}")
    #print(f"affiliation:{affiliations}")

    return total

def prepare_data_v0():
    total = 0
    read_xlsx = load_workbook(r'list.xlsx')
    read_sheet = read_xlsx.active
    name_col = read_sheet['B']
    affiliation_col = read_sheet['C']

    for cell in name_col:
        names.append(cell.value)

    for cell in affiliation_col:
        affiliations.append(cell.value)
        total+=1

    total-=1
    del names[0]
    del affiliations[0]

    return total

## This is master template slide
prs = Presentation('template.pptx')

def copy_slide(prs, index):
    template = prs.slides[index]
    try:
        blank_slide_layout = prs.slide_layouts.get_by_name()
    except:
        blank_slide_layout = prs.slide_layouts[6]

    copied_slide = prs.slides.add_slide(blank_slide_layout)
    
    ## (jhpark): I refer to this [link](https://stackoverflow.com/a/62921781)
    imgDict = {}
    for shape in template.shapes:
        elem = shape.element
        if 'logo' in shape.name:
            with open(shape.name+'.jpg', 'wb') as f:
                f.write(shape.image.blob)
            imgDict[shape.name+'.jpg'] = [shape.left, shape.top, shape.width, shape.height]
        else:        
            new_elem = copy.deepcopy(elem)
            copied_slide.shapes._spTree.insert_element_before(new_elem, 'p:extLst')        
    
    for k, v in imgDict.items():
        copied_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)
    return copied_slide

# 4개씩 이름, 소속 끊어서 저장
def write_slide(prs, index, names, affiliations):
    cur = prs.slides[index]
    name_idx=0
    affiliation_idx=0

    for shape in cur.shapes:
        if 'name' in shape.name:
            tf=shape.text_frame
            tf.clear()
            para=tf.paragraphs[0]
            para.alignment=PP_ALIGN.CENTER
            run=para.add_run()
            if name_idx < len(names):
                run.text=names[name_idx]
            else:
                run.text=''
            font=run.font
            font.size=Pt(36)
            font.name="Malgun Gothic"
            name_idx+=1

        if 'affiliation' in shape.name:
            tf=shape.text_frame
            tf.clear()
            para=tf.paragraphs[0]
            para.alignment=PP_ALIGN.CENTER
            run=para.add_run()
            if affiliation_idx < len(affiliations):
                run.text=affiliations[affiliation_idx]
            else:
                run.text=''
            font=run.font
            font.size=Pt(20)
            font.name="Malgun Gothic"
            affiliation_idx+=1

## (jhpark): I refer to this [link](https://stackoverflow.com/questions/50866634/how-to-copy-a-slide-with-python-pptx)
def delete_slide(prs, index):
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = prs.slides[index].slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

def make(prs,total,names,affiliations):
    ## step1. calculate total # of iterations
    ## step2. Names and affiliations are retrieved in groups of *four* to create name tags.

    slide_num=1
    iterations=math.ceil(total/4)
    for i in range(iterations):
        copy_slide(prs, 0)

    for i in range(0, total, 4):
        cur_name = names[i:i+4]
        cur_affiliation =affiliations[i:i+4]
        write_slide(prs,slide_num,cur_name,cur_affiliation)
        slide_num+=1

    # delete master template slide
    delete_slide(prs,0)
    prs.save('output.pptx')
      

def main():
    #print("Welcom Nametag v0.2")
    #total = prepare_data_v1()

    print("Welcom Nametag v0.1")
    total = prepare_data_v0()
    print(f"Total # of registants: {total}")
    make(prs,total,names,affiliations)


if __name__ == "__main__":
    main()